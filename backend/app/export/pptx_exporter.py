"""PPTX export for persisted quizzes — quiz-show presentation style."""

from __future__ import annotations

import random
from io import BytesIO

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches
from pptx.util import Pt

from backend.app.domain.errors import DomainValidationError
from backend.app.domain.models import Question
from backend.app.domain.models import Quiz
from backend.app.export.base import ExportedQuizFile

_SLIDE_W = Inches(13.33)
_SLIDE_H = Inches(7.5)

_BG           = RGBColor(0x08, 0x10, 0x28)
_QUESTION_BOX = RGBColor(0x0A, 0x0A, 0x0A)
_OPTION_BOX   = RGBColor(0x0D, 0x0D, 0x1A)
_GRAD_1       = RGBColor(0x10, 0x30, 0xD0)
_GRAD_2       = RGBColor(0x99, 0x11, 0xEE)
_GLOW_CLR     = RGBColor(0x60, 0xA0, 0xFF)
_ARROW_CLR    = RGBColor(0x88, 0xCC, 0xFF)
_WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
_PURPLE       = RGBColor(0xBB, 0x44, 0xFF)
_BADGE_BG     = RGBColor(0x22, 0x2A, 0x44)
_BORDER       = RGBColor(0x33, 0x44, 0x66)
_HINT         = RGBColor(0x55, 0x66, 0x88)

_LABELS = ["A", "B", "C", "D"]


def _xml_set_bg(slide, color: RGBColor) -> None:
    spTree = slide.shapes._spTree
    bg_xml = (
        '<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        "<p:bgPr>"
        f'<a:solidFill><a:srgbClr val="{str(color)}"/></a:solidFill>'
        "<a:effectLst/>"
        "</p:bgPr>"
        "</p:bg>"
    )
    bg_elem = etree.fromstring(bg_xml)
    spTree.getparent().insert(list(spTree.getparent()).index(spTree), bg_elem)


def _xml_add_glow(shape, color: RGBColor, radius_pt: float = 12.0) -> None:
    radius_emu = int(radius_pt * 12700)
    sp = shape._element
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        return
    for existing in spPr.findall(qn("a:effectLst")):
        spPr.remove(existing)
    spPr.append(etree.fromstring(
        f'<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:glow rad="{radius_emu}">'
        f'<a:srgbClr val="{str(color)}"><a:alpha val="70000"/></a:srgbClr>'
        f"</a:glow></a:effectLst>"
    ))


def _xml_grad_fill(shape, c1: RGBColor, c2: RGBColor) -> None:
    sp = shape._element
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        return
    for tag in ("a:solidFill", "a:gradFill", "a:noFill", "a:pattFill"):
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)
    grad = etree.fromstring(
        '<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        "<a:gsLst>"
        f'<a:gs pos="0"><a:srgbClr val="{str(c1)}"/></a:gs>'
        f'<a:gs pos="100000"><a:srgbClr val="{str(c2)}"/></a:gs>'
        "</a:gsLst>"
        '<a:lin ang="5400000" scaled="0"/>'
        "</a:gradFill>"
    )
    prstGeom = spPr.find(qn("a:prstGeom"))
    pos = list(spPr).index(prstGeom) + 1 if prstGeom is not None else 1
    spPr.insert(pos, grad)


def _add_hexagon(slide, left, top, width, height, fill: RGBColor, border: RGBColor, border_pt: float = 1.2):
    shape = slide.shapes.add_shape(6, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(border_pt)
    return shape


def _add_hexagon_gradient(slide, left, top, width, height, c1: RGBColor, c2: RGBColor, border: RGBColor, border_pt: float = 3.0):
    shape = slide.shapes.add_shape(6, left, top, width, height)
    _xml_grad_fill(shape, c1, c2)
    shape.line.color.rgb = border
    shape.line.width = Pt(border_pt)
    return shape


def _add_rounded_rect(slide, left, top, width, height, fill: RGBColor, border: RGBColor, border_pt: float = 1.0):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(border_pt)
    return shape


def _textbox(slide, left, top, width, height, text: str, size_pt: float, color: RGBColor,
             bold: bool = False, align=PP_ALIGN.LEFT, v_anchor=MSO_ANCHOR.MIDDLE) -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = v_anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.color.rgb = color
    run.font.bold = bold


def _make_slide(prs: Presentation, q_index: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _xml_set_bg(slide, _BG)
    badge = _add_rounded_rect(slide, Inches(0.35), Inches(0.22), Inches(1.6), Inches(0.38), _BADGE_BG, _BORDER)
    tf = badge.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"Вопрос {q_index}"
    run.font.size = Pt(13)
    run.font.color.rgb = _WHITE
    run.font.bold = True
    return slide


def _add_question_box(slide, prompt: str, top=Inches(0.9), height=Inches(1.9)) -> None:
    shape = _add_hexagon(slide, Inches(1.5), top, Inches(10.3), height, _QUESTION_BOX, _BORDER, 1.2)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = prompt
    run.font.size = Pt(20)
    run.font.color.rgb = _WHITE
    run.font.bold = True


def _add_option_pill(slide, ox, oy, w, h, label: str, text: str, highlighted: bool) -> None:
    lw = Inches(0.65)
    if highlighted:
        sh = _add_hexagon_gradient(slide, ox, oy, w, h, _GRAD_1, _GRAD_2, _GLOW_CLR, 3.0)
        _xml_add_glow(sh, _GLOW_CLR, 16)
    else:
        _add_hexagon(slide, ox, oy, w, h, _OPTION_BOX, _BORDER, 1.2)
    _textbox(slide, ox + Inches(0.15), oy, lw, h, label, 20, _PURPLE, bold=True, align=PP_ALIGN.CENTER)
    _textbox(slide, ox + lw + Inches(0.2), oy, w - lw - Inches(0.35), h, text, 17, _WHITE)


def _add_answer_reveal(slide, answer: str) -> None:
    box = _add_rounded_rect(slide, Inches(2.0), Inches(3.8), Inches(9.3), Inches(1.5),
                            _GRAD_1, _GLOW_CLR, 2.5)
    _xml_add_glow(box, _GLOW_CLR, 18)
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = answer
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = _WHITE


def _add_arrow(slide, x1, y1, x2, y2, color: RGBColor, width_pt: float = 2.5) -> None:
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = Pt(width_pt)
    sp = connector._element
    ln = sp.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}ln")
    if ln is not None:
        ln.append(etree.fromstring(
            '<a:tailEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
            ' type="arrow" w="med" len="med"/>'
        ))


class QuizPptxExporter:
    """Export persisted quizzes into quiz-show style PPTX presentations.

    Each question produces two slides:
    - Slide A: question with options, no highlight (shown to audience)
    - Slide B: same question with correct answer revealed
    """

    media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    def export(self, quiz: Quiz) -> ExportedQuizFile:
        """Render one quiz into a quiz-show PPTX file."""

        prs = Presentation()
        prs.slide_width = _SLIDE_W
        prs.slide_height = _SLIDE_H
        for index, question in enumerate(quiz.questions, start=1):
            self._add_question_slides(prs, index, question)
        output = BytesIO()
        prs.save(output)
        return ExportedQuizFile(
            filename=f"{quiz.quiz_id}.pptx",
            media_type=self.media_type,
            content_bytes=output.getvalue(),
        )

    def _add_question_slides(self, prs: Presentation, index: int, question: Question) -> None:
        qt = question.question_type
        if qt in {"single_choice", "true_false"}:
            self._slides_choice(prs, index, question)
        elif qt in {"fill_blank", "short_answer"}:
            self._slides_open(prs, index, question)
        elif qt == "matching":
            self._slides_matching(prs, index, question)
        else:
            self._slides_choice(prs, index, question)

    def _slides_choice(self, prs: Presentation, index: int, question: Question) -> None:
        options = list(question.options)
        correct_idx = self._resolve_correct_option_index(question)
        is_two = question.question_type == "true_false" or len(options) <= 2

        if is_two:
            col_left = [Inches(1.2), Inches(7.1)]
            row_top  = [Inches(4.0)]
            opt_w, opt_h = Inches(5.0), Inches(1.3)
        else:
            col_left = [Inches(0.6), Inches(6.8)]
            row_top  = [Inches(3.2), Inches(4.65)]
            opt_w, opt_h = Inches(5.9), Inches(1.1)

        for reveal in (False, True):
            slide = _make_slide(prs, index)
            _add_question_box(slide, question.prompt)
            for i, option in enumerate(options[:4]):
                col = i % 2
                row = i // 2
                ox = col_left[col]
                oy = row_top[row] if row < len(row_top) else row_top[-1] + (row - len(row_top) + 1) * (opt_h + Inches(0.15))
                _add_option_pill(slide, ox, oy, opt_w, opt_h,
                                 _LABELS[i], option.text,
                                 highlighted=(reveal and i == correct_idx))

    def _slides_open(self, prs: Presentation, index: int, question: Question) -> None:
        answer = question.correct_answer or ""
        slide_a = _make_slide(prs, index)
        _add_question_box(slide_a, question.prompt, top=Inches(2.8), height=Inches(2.0))
        _textbox(slide_a, Inches(1.5), Inches(5.2), Inches(10.3), Inches(0.6),
                 "Ответьте устно…", 16, _HINT, align=PP_ALIGN.CENTER)
        slide_b = _make_slide(prs, index)
        _add_question_box(slide_b, question.prompt, top=Inches(1.5), height=Inches(1.9))
        _add_answer_reveal(slide_b, answer)

    def _slides_matching(self, prs: Presentation, index: int, question: Question) -> None:
        pairs = list(question.matching_pairs)
        lefts  = [p.left  for p in pairs]
        rights = [p.right for p in pairs]
        shuffled = rights[:]
        for _ in range(20):
            random.shuffle(shuffled)
            if shuffled != rights:
                break

        col_w   = Inches(5.5)
        row_h   = Inches(0.9)
        gap     = Inches(0.15)
        left_x  = Inches(0.5)
        right_x = Inches(7.3)
        start_y = Inches(2.5)

        for reveal in (False, True):
            slide = _make_slide(prs, index)
            _add_question_box(slide, question.prompt, top=Inches(0.85), height=Inches(1.4))
            right_col = rights if reveal else shuffled
            right_text_to_y: dict[str, int] = {}

            for row, (ltext, rtext) in enumerate(zip(lefts, right_col)):
                oy = start_y + row * (row_h + gap)
                _add_option_pill(slide, left_x, oy, col_w, row_h, _LABELS[row], ltext, highlighted=False)
                _add_hexagon(slide, right_x, oy, col_w, row_h, _OPTION_BOX, _BORDER, 1.2)
                _textbox(slide, right_x + Inches(0.2), oy, col_w - Inches(0.25), row_h, rtext, 16, _WHITE)
                right_text_to_y[rtext] = oy

            if reveal:
                for row, pair in enumerate(pairs):
                    ly = start_y + row * (row_h + gap)
                    ry = right_text_to_y.get(pair.right)
                    if ry is None:
                        continue
                    _add_arrow(slide,
                               left_x + col_w, ly + row_h // 2,
                               right_x,        ry + row_h // 2,
                               _ARROW_CLR, 2.5)

    @staticmethod
    def _resolve_correct_option_index(question: Question) -> int:
        if (
            question.correct_option_index is None
            or question.correct_option_index < 0
            or question.correct_option_index >= len(question.options)
        ):
            raise DomainValidationError(
                f"question '{question.question_id}' has invalid correct_option_index for PPTX export"
            )
        return question.correct_option_index
