from __future__ import annotations

from dataclasses import replace
from io import BytesIO

import pytest
from pptx import Presentation

from backend.app.domain.errors import DomainValidationError
from backend.app.domain.models import Explanation
from backend.app.domain.models import Option
from backend.app.domain.models import Question
from backend.app.domain.models import Quiz
from backend.app.export.pptx_exporter import QuizPptxExporter
from backend.app.export.registry import DEFAULT_QUIZ_EXPORT_REGISTRY


def build_quiz() -> Quiz:
    return Quiz(
        quiz_id="quiz-ru-1",
        document_id="doc-ru-1",
        title="Тренировочный квиз по географии",
        version=3,
        last_edited_at="2026-04-22T12:00:00.000000Z",
        questions=(
            Question(
                question_id="question-1",
                prompt="Какой город является столицей России?",
                options=(
                    Option(option_id="option-1", text="Москва"),
                    Option(option_id="option-2", text="Казань"),
                ),
                correct_option_index=0,
                explanation=Explanation(text="Москва является столицей России."),
            ),
            Question(
                question_id="question-2",
                prompt="Какая река протекает через Санкт-Петербург?",
                options=(
                    Option(option_id="option-1", text="Волга"),
                    Option(option_id="option-2", text="Нева"),
                ),
                correct_option_index=1,
                explanation=None,
            ),
        ),
    )


def read_pptx_text(content_bytes: bytes) -> tuple[str, int]:
    presentation = Presentation(BytesIO(content_bytes))
    slide_text: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            slide_text.append(shape.text)
    return "\n".join(slide_text), len(presentation.slides)


def test_pptx_exporter_builds_openable_quiz_presentation_with_cyrillic_content() -> None:
    exporter = QuizPptxExporter()

    exported_file = exporter.export(build_quiz())
    presentation_text, slide_count = read_pptx_text(exported_file.content_bytes)

    assert exported_file.filename == "quiz-ru-1.pptx"
    assert exported_file.media_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    assert exported_file.content_bytes.startswith(b"PK")
    assert slide_count == 2
    assert "Тренировочный квиз по географии" in presentation_text
    assert "Какой город является столицей России?" in presentation_text
    assert "1. Москва" in presentation_text
    assert "2. Казань" in presentation_text
    assert "Правильный ответ: Москва" in presentation_text
    assert "Пояснение: Москва является столицей России." in presentation_text
    assert "Какая река протекает через Санкт-Петербург?" in presentation_text
    assert "Правильный ответ: Нева" in presentation_text


def test_default_export_registry_exposes_pptx_exporter() -> None:
    exported_file = DEFAULT_QUIZ_EXPORT_REGISTRY.export(build_quiz(), "PPTX")
    presentation_text, slide_count = read_pptx_text(exported_file.content_bytes)

    assert "pptx" in DEFAULT_QUIZ_EXPORT_REGISTRY.supported_formats()
    assert exported_file.filename == "quiz-ru-1.pptx"
    assert slide_count == 2
    assert "Тренировочный квиз по географии" in presentation_text


def test_pptx_exporter_rejects_invalid_correct_option_index() -> None:
    quiz = build_quiz()
    invalid_question = replace(quiz.questions[0], correct_option_index=5)
    invalid_quiz = replace(quiz, questions=(invalid_question,))

    with pytest.raises(DomainValidationError, match="correct_option_index"):
        QuizPptxExporter().export(invalid_quiz)
