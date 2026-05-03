from __future__ import annotations

from io import BytesIO

from docx import Document
from fastapi.testclient import TestClient
from pptx import Presentation

from backend.app.core.config import AppConfig
from backend.app.domain.models import Explanation
from backend.app.domain.models import Option
from backend.app.domain.models import ProviderHealthStatus
from backend.app.domain.models import Question
from backend.app.domain.models import Quiz
from backend.app.main import create_app
from backend.app.storage.quizzes import FileSystemQuizRepository


class StubProvider:
    """Provider test double for advanced export API flows."""

    def healthcheck(self) -> ProviderHealthStatus:
        return ProviderHealthStatus(status="available", message="LM Studio is available")

    def generate_structured(self, request):
        raise AssertionError("generate_structured should not be called by export tests")

    def embed(self, request):
        raise AssertionError("embed should not be called by export tests")


def build_config() -> AppConfig:
    return AppConfig(
        lm_studio_base_url="http://localhost:1234/v1",
        lm_studio_model="local-model",
        log_format="%(levelname)s:%(message)s",
    )


def build_quiz() -> Quiz:
    return Quiz(
        quiz_id="quiz-ru-1",
        document_id="doc-ru-1",
        title="Тренировочный квиз по географии",
        version=0,
        last_edited_at="",
        questions=(
            Question(
                question_id="question-1",
                prompt="Какой город является столицей России?",
                options=(
                    Option(option_id="option-1", text="Москва"),
                    Option(option_id="option-2", text="Казань"),
                ),
                correct_option_index=0,
                explanation=Explanation(text="Москва является столицей России."),
            ),
        ),
    )


def create_test_client(tmp_path) -> TestClient:
    app = create_app(config=build_config(), provider=StubProvider(), storage_root=tmp_path)
    return TestClient(app)


def save_quiz(tmp_path) -> Quiz:
    repository = FileSystemQuizRepository(tmp_path)
    return repository.save(build_quiz())


def read_docx_text(content_bytes: bytes) -> str:
    document = Document(BytesIO(content_bytes))
    return "\n".join(paragraph.text for paragraph in document.paragraphs)


def read_pptx_text(content_bytes: bytes) -> str:
    presentation = Presentation(BytesIO(content_bytes))
    slide_text: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            slide_text.append(shape.text)
    return "\n".join(slide_text)


def test_export_formats_endpoint_reports_supported_formats(tmp_path) -> None:
    client = create_test_client(tmp_path)

    response = client.get("/export/formats")

    assert response.status_code == 200
    assert response.json()["request_id"] == response.headers["X-Request-ID"]
    formats = {item["format"]: item["media_type"] for item in response.json()["formats"]}
    assert formats == {
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "json": "application/json; charset=utf-8",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }


def test_docx_export_endpoint_downloads_openable_cyrillic_quiz_file(tmp_path) -> None:
    persisted_quiz = save_quiz(tmp_path)
    client = create_test_client(tmp_path)

    response = client.get(f"/quizzes/{persisted_quiz.quiz_id}/export/docx")

    assert response.status_code == 200
    assert response.headers["X-Request-ID"]
    assert response.headers["content-type"] == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    assert response.headers["content-disposition"] == 'attachment; filename="quiz-ru-1.docx"'
    assert response.content.startswith(b"PK")
    document_text = read_docx_text(response.content)
    assert "Тренировочный квиз по географии" in document_text
    assert "Какой город является столицей России?" in document_text
    assert "Правильный ответ: Москва" in document_text


def test_pptx_export_endpoint_downloads_openable_cyrillic_quiz_file(tmp_path) -> None:
    persisted_quiz = save_quiz(tmp_path)
    client = create_test_client(tmp_path)

    response = client.get(f"/quizzes/{persisted_quiz.quiz_id}/export/pptx")

    assert response.status_code == 200
    assert response.headers["X-Request-ID"]
    assert response.headers["content-type"] == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    assert response.headers["content-disposition"] == 'attachment; filename="quiz-ru-1.pptx"'
    assert response.content.startswith(b"PK")
    presentation_text = read_pptx_text(response.content)
    assert "Какой город является столицей России?" in presentation_text
    assert "Москва" in presentation_text


def test_advanced_export_endpoint_rejects_unsupported_format(tmp_path) -> None:
    persisted_quiz = save_quiz(tmp_path)
    client = create_test_client(tmp_path)

    response = client.get(f"/quizzes/{persisted_quiz.quiz_id}/export/pdf")

    assert response.status_code == 400
    assert response.json()["error"]["code"] == "unsupported_export_format"
    assert "supported formats: docx, json, pptx" in response.json()["error"]["message"]
    assert response.json()["request_id"] == response.headers["X-Request-ID"]
